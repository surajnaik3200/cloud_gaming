from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

BASE_DIR = os.path.join(os.path.dirname(__file__), '..')
IMAGES_DIR = os.path.join(BASE_DIR, 'assets', 'images')
OUT_PATH = os.path.join(BASE_DIR, 'presentation.pptx')
LOGO_PATH = os.path.join(IMAGES_DIR, 'persius_logo.png')

# Theme colors from site's CSS
THEME = {
    'bg': (0x0a, 0x0e, 0x27),          # --bg
    'primary': (0xb3, 0x00, 0xff),     # --primary
    'primary_light': (0xff, 0x00, 0xff),
    'text': (0xe0, 0xe0, 0xe0),        # --text
    'text_bright': (0x00, 0xff, 0xff), # --text-bright
}

headings = [
    ('Admin Panel', [
        'Central management interface',
        'Quick access to site sections',
        'User and content controls'
    ]),
    ('Dashboard', [
        'Overview metrics and KPIs',
        'Recent activity and alerts',
        'Quick links to sections'
    ]),
    ('Games', [
        'Add, edit, delete game entries',
        'List with filters and search',
        'Manage game assets and posters'
    ]),
    ('Players', [
        'View and manage player accounts',
        'Search and filter players',
        'Assign roles and permissions'
    ]),
    ('Subscriptions', [
        'Manage subscription plans',
        'View active subscribers',
        'Handle renewals and cancellations'
    ]),
    ('Payments', [
        'List and reconcile payments',
        'Process refunds and disputes'
    ]),
    ('Logout', [
        'Securely sign out admin session',
        'Clear session and auth tokens'
    ])
]


def set_title_style(shape):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Arial'
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = RGBColor(*THEME['primary_light'])


def set_bullet_style(shape):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Arial'
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(*THEME['text'])


def apply_theme_to_slide(slide):
    # Set solid background color
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*THEME['bg'])
    except Exception:
        pass
    # Add logo if present
    try:
        if os.path.isfile(LOGO_PATH):
            slide.shapes.add_picture(LOGO_PATH, Inches(0.3), Inches(0.25), height=Inches(0.6))
    except Exception:
        pass


def create_presentation():
    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = 'Admin Panel Overview'
    subtitle.text = 'Generated from project sidebar — 2–3 points per section'
    set_title_style(title)
    set_bullet_style(subtitle)
    apply_theme_to_slide(slide)

    # Content slides
    for heading, bullets in headings:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        title = slide.shapes.title
        body = slide.shapes.placeholders[1]
        title.text = heading
        p = body.text_frame
        p.clear()
        for i, b in enumerate(bullets):
            if i == 0:
                p.text = b
            else:
                p.add_paragraph().text = b
        set_title_style(title)
        set_bullet_style(body)
        apply_theme_to_slide(slide)

    # Image slides: one image per slide
    if os.path.isdir(IMAGES_DIR):
        images = [f for f in os.listdir(IMAGES_DIR) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif'))]
        for img in images:
            img_path = os.path.join(IMAGES_DIR, img)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            apply_theme_to_slide(slide)
            # place image centered-ish with safe sizing
            left = Inches(1)
            top = Inches(1)
            height = Inches(5.5)
            try:
                slide.shapes.add_picture(img_path, left, top, height=height)
            except Exception:
                continue

    prs.save(OUT_PATH)
    print(f'Created presentation at: {OUT_PATH}')


if __name__ == '__main__':
    create_presentation()
